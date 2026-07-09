"""Generate the NVIDIA-meeting slide deck for the ICICLE Edge Fleet Control Plane.

Run:  python deploy/make_nvidia_deck.py
Out:  deploy/ICICLE_Edge_Platform_NVIDIA.pptx

Pure python-pptx, no template — a dark 16:9 deck themed to match the portal
(indigo/violet brand on deep-navy), tuned for a technical Jetson/NVIDIA audience.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- Palette (matches the portal: indigo/violet brand on deep-navy surface) --
ACCENT = RGBColor(0x81, 0x8C, 0xF8)   # brand-400 (reads bright on dark navy)
ACCENT_D = RGBColor(0x4F, 0x46, 0xE5)  # brand-600
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)   # secondary brand accent
BG = RGBColor(0x0B, 0x10, 0x20)        # portal sidebar navy (--side-bg)
PANEL = RGBColor(0x13, 0x1A, 0x2E)     # --side-bg-soft
PANEL_L = RGBColor(0x1E, 0x27, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9F, 0xAA, 0xC4)     # muted slate-indigo
LINE = RGBColor(0x2A, 0x33, 0x50)

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

_page = 0


def _set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _txt(slide, left, top, width, height, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def _run(p, text, *, size=18, color=WHITE, bold=False, italic=False, font="Segoe UI"):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def _footer(slide):
    global _page
    _page += 1
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.18), EMU_W, Pt(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LINE
    bar.line.fill.background()
    _, tf = _txt(slide, Inches(0.5), Inches(7.0), Inches(8), Inches(0.4))
    _run(tf.paragraphs[0], "ICICLE Edge Fleet Control Plane", size=10, color=MUTED)
    tb2, tf2 = _txt(slide, Inches(11.5), Inches(7.0), Inches(1.33), Inches(0.4))
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _run(p, f"{_page:02d}", size=10, color=ACCENT, bold=True)


def content_slide(kicker, title):
    slide = prs.slides.add_slide(BLANK)
    _set_bg(slide)
    # kicker
    _, ktf = _txt(slide, Inches(0.7), Inches(0.45), Inches(11), Inches(0.4))
    _run(ktf.paragraphs[0], kicker.upper(), size=13, color=ACCENT, bold=True)
    # title
    _, ttf = _txt(slide, Inches(0.7), Inches(0.8), Inches(12), Inches(1.0))
    _run(ttf.paragraphs[0], title, size=32, color=WHITE, bold=True)
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.62), Inches(1.6), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    _footer(slide)
    return slide


def bullets(slide, items, *, left=Inches(0.8), top=Inches(2.0),
            width=Inches(11.7), height=Inches(4.6), size=18, gap=10):
    _, tf = _txt(slide, left, top, width, height)
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        if level == 0:
            _run(p, "▍ ", size=size, color=ACCENT, bold=True)
            _run(p, text, size=size, color=WHITE)
        elif level == 1:
            p.level = 1
            _run(p, "—  ", size=size - 2, color=ACCENT)
            _run(p, text, size=size - 2, color=MUTED)
        else:  # note line
            _run(p, text, size=size - 3, color=MUTED, italic=True)
    return tf


def box(slide, left, top, width, height, title, sub=None, *,
        fill=PANEL, border=ACCENT, title_color=WHITE, sub_color=MUTED, tsize=15, ssize=11):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = border
    sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, title, size=tsize, color=title_color, bold=True)
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        _run(p2, sub, size=ssize, color=sub_color)
    return sp


def arrow(slide, left, top, width, height=Inches(0.4), color=ACCENT, shape=MSO_SHAPE.RIGHT_ARROW):
    sp = slide.shapes.add_shape(shape, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


# ============================================================================
# SLIDE 1 — TITLE
# ============================================================================
s = prs.slides.add_slide(BLANK)
_set_bg(s)
# big green side band
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), EMU_H)
band.fill.solid(); band.fill.fore_color.rgb = ACCENT; band.line.fill.background()
_, tf = _txt(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(0.5))
_run(tf.paragraphs[0], "NSF ICICLE AI INSTITUTE  ·  EDGE TOOLING", size=15, color=ACCENT, bold=True)
_, tf = _txt(s, Inches(0.9), Inches(2.6), Inches(11.6), Inches(1.8))
_run(tf.paragraphs[0], "ICICLE Edge Fleet Control Plane", size=48, color=WHITE, bold=True)
_, tf = _txt(s, Inches(0.92), Inches(4.05), Inches(11.4), Inches(1.0))
_run(tf.paragraphs[0],
     "Fleet management and AI model deployment for NVIDIA Jetson at the edge",
     size=22, color=MUTED)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(5.0), Inches(3.2), Pt(3))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
_, tf = _txt(s, Inches(0.92), Inches(5.25), Inches(11), Inches(0.8))
_run(tf.paragraphs[0], "Prepared for NVIDIA  ·  Technical overview & collaboration", size=16, color=WHITE)
p = tf.add_paragraph(); _run(p, "Live demo available on request", size=13, color=MUTED)

# ============================================================================
# SLIDE 2 — THE CHALLENGE
# ============================================================================
s = content_slide("The problem", "Running AI on a fleet of edge devices is hard")
bullets(s, [
    "AI increasingly runs in the field — on cameras, sensors, robots — not just in the cloud.",
    ("Each node is an NVIDIA Jetson: a GPU computer on a farm, in a lab, or in the wild.", 1),
    "Getting a model onto every device meant SSH-ing in, copying files, and running commands by hand.",
    ("Slow, error-prone, and impossible to track across dozens of devices.", 1),
    "No single view of what's deployed where, whether it's healthy, or who changed it.",
    "Field devices sit behind NAT/firewalls — you usually can't just reach into them.",
], top=Inches(2.05))

# ============================================================================
# SLIDE 3 — WHAT WE BUILT
# ============================================================================
s = content_slide("The platform", "One control plane for the whole edge fleet")
bullets(s, [
    "A web control plane to publish AI \u201cmodel cards\u201d (a model + exactly how to run it).",
    "Deploy to one device, a group, or an entire device generation — in a few clicks.",
    "Monitor, stop, restart, and remove deployments live — no terminal required.",
    "A lightweight agent on each Jetson does the heavy lifting and reports back.",
    "Cloud-native: runs as a container on Tapis Pods (Kubernetes); integrates with the ICICLE ecosystem.",
    "Live status everywhere — Delivering, Running, Stopped — with a full audit trail.",
], top=Inches(2.05))

# ============================================================================
# SLIDE 4 — ARCHITECTURE DIAGRAM
# ============================================================================
s = content_slide("Architecture", "How the pieces fit together")
top = Inches(2.25)
bw, bh = Inches(3.5), Inches(1.5)
box(s, Inches(0.7), top, bw, bh, "Operator", "Browser dashboard", fill=PANEL_L)
box(s, Inches(4.9), top, bw, bh, "Control Plane", "FastAPI on Tapis Pods", border=ACCENT)
box(s, Inches(9.1), top, bw, bh, "NVIDIA Jetson Fleet", "Edge agent + Docker", border=ACCENT)
arrow(s, Inches(4.25), Inches(2.75), Inches(0.6))
# bidirectional between control plane and fleet
arrow(s, Inches(8.45), Inches(2.62), Inches(0.6))
arrow(s, Inches(8.45), Inches(3.0), Inches(0.6), color=MUTED, shape=MSO_SHAPE.LEFT_ARROW)
_, tf = _txt(s, Inches(8.2), Inches(3.45), Inches(2.0), Inches(0.3))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_run(p, "MQTT cmds  /  HTTPS heartbeats", size=9, color=MUTED)
# ecosystem row
ey = Inches(4.55)
ew = Inches(2.85)
box(s, Inches(0.7), ey, ew, Inches(1.25), "PostgreSQL", "Devices, models, audit", fill=PANEL, border=LINE)
box(s, Inches(3.75), ey, ew, Inches(1.25), "Patra", "Model provenance", fill=PANEL, border=LINE)
box(s, Inches(6.8), ey, ew, Inches(1.25), "MQTT Broker", "Command delivery", fill=PANEL, border=LINE)
box(s, Inches(9.85), ey, ew, Inches(1.25), "Tapis", "Auth (OAuth2) + Pods hosting", fill=PANEL, border=LINE)
_, tf = _txt(s, Inches(0.7), Inches(4.2), Inches(11), Inches(0.3))
_run(tf.paragraphs[0], "Shared ICICLE / Tapis infrastructure", size=12, color=ACCENT, bold=True)

# ============================================================================
# SLIDE 5 — THE JETSON AGENT
# ============================================================================
s = content_slide("On the device", "The NVIDIA Jetson agent")
bullets(s, [
    "A small Python systemd service installed on each Jetson via a generated installer.",
    "Subscribes to MQTT command topics (per-device, per-group, per-generation).",
    "Runs models as Docker containers on the device — deploy / stop / restart / delete.",
    "Reports heartbeats over HTTPS: status, running containers, and GPU telemetry.",
    "Makes an outbound connection only — works behind NAT/firewalls with no inbound ports.",
    "Handles real-world Jetson needs (e.g. X11/display + camera access for CV containers).",
], top=Inches(2.05))

# ============================================================================
# SLIDE 6 — MODEL LIFECYCLE
# ============================================================================
s = content_slide("Model lifecycle", "From model card to a running container on the GPU")
bullets(s, [
    "A model card is self-contained: artifact (e.g. a TensorRT .engine), image, env, mounts, args.",
    "Deploy builds one MQTT payload — the agent needs no extra lookups.",
    "Agent resolves & validates the artifact, pulls the image, and launches the container.",
    ("docker run with --runtime nvidia / --gpus, privileged & device mounts as the card specifies.", 1),
    "Live status flows back: Delivering \u2192 Pulling \u2192 Starting \u2192 Running.",
    "Stop keeps image + artifacts cached for instant restart; Delete fully cleans the device.",
], top=Inches(2.05))

# ============================================================================
# SLIDE 7 — DEPLOYMENT FLOW DIAGRAM
# ============================================================================
s = content_slide("Deployment flow", "One click → a container on the Jetson")
fy = Inches(2.5)
fbw, fbh = Inches(2.5), Inches(1.6)
gaps = Inches(0.45)
xs = [Inches(0.55), Inches(3.5), Inches(6.45), Inches(9.4)]
labels = [
    ("Deploy", "Operator picks a model + target"),
    ("Build payload", "Control plane assembles model card"),
    ("MQTT command", "Delivered to the device topic"),
    ("docker run", "Agent starts the GPU container"),
]
for x, (t, sub) in zip(xs, labels):
    box(s, x, fy, fbw, fbh, t, sub, border=ACCENT if t in ("MQTT command", "docker run") else LINE)
for x in xs[:-1]:
    arrow(s, x + fbw + Emu(int(Inches(0.02))), Inches(3.1), Inches(0.4))
# return path
ret = box(s, Inches(3.5), Inches(4.7), Inches(8.4), Inches(1.1),
          "ACK + heartbeats", "RUNNING / FAILED status and GPU health flow back over HTTPS — dashboard updates live",
          fill=PANEL, border=ACCENT, ssize=12)
arrow(s, Inches(10.0), Inches(4.25), Inches(0.4), Inches(0.45),
      color=MUTED, shape=MSO_SHAPE.DOWN_ARROW)

# ============================================================================
# SLIDE 8 — LEVERAGING THE NVIDIA STACK  (key slide)
# ============================================================================
s = content_slide("Why it matters to NVIDIA", "Built on the NVIDIA edge stack")
bullets(s, [
    "Jetson-native: containers run with the NVIDIA container runtime and GPU access.",
    "TensorRT engines: model cards ship & validate .engine artifacts for accelerated inference.",
    "Any CUDA / TensorRT container image runs unmodified — bring your own model image.",
    "GPU observability: per-device telemetry (GPU temp, memory, running containers).",
    "DeepStream-ready: the container model fits DeepStream / Triton pipelines as a next step.",
    "Scales across the Jetson family — Nano today, Orin / IGX tomorrow.",
], top=Inches(1.95), size=18)

# ============================================================================
# SLIDE — MULTI-CAMERA HARDWARE VALIDATION (figure)
# ============================================================================
s = content_slide("Validated on hardware", "Multi-camera capture on a single Jetson")

# --- Left: two active IP cameras ---
cam_x = Inches(0.6)
box(s, cam_x, Inches(2.55), Inches(2.25), Inches(1.0), "IP Camera 1", "PoE-powered", border=ACCENT)
box(s, cam_x, Inches(3.95), Inches(2.25), Inches(1.0), "IP Camera 2", "PoE-powered", border=ACCENT)

# PoE link arrows + label
arrow(s, Inches(2.95), Inches(2.92), Inches(0.65), Inches(0.28))
arrow(s, Inches(2.95), Inches(4.32), Inches(0.65), Inches(0.28))
_, tf = _txt(s, Inches(2.78), Inches(3.55), Inches(1.0), Inches(0.6))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_run(p, "PoE", size=11, color=MUTED, bold=True)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
_run(p, "Ethernet", size=11, color=MUTED)

# --- Middle: the Jetson board ---
jx, jy, jw, jh = Inches(3.75), Inches(2.25), Inches(4.55), Inches(3.4)
jet = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, jx, jy, jw, jh)
jet.fill.solid(); jet.fill.fore_color.rgb = PANEL_L
jet.line.color.rgb = ACCENT; jet.line.width = Pt(1.5); jet.shadow.inherit = False
_, tf = _txt(s, jx, Inches(2.4), jw, Inches(0.45))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_run(p, "NVIDIA Jetson", size=17, color=WHITE, bold=True)

# 4 PoE ports (2 in use, 2 free)
port_w, port_h = Inches(0.62), Inches(0.46)
port_tops = Inches(3.0)
for idx in range(4):
    px = Emu(int(jx) + int(Inches(0.28)) + idx * int(Inches(0.78)))
    used = idx < 2
    pr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, port_tops, port_w, port_h)
    pr.fill.solid(); pr.fill.fore_color.rgb = ACCENT if used else BG
    pr.line.color.rgb = ACCENT if used else LINE; pr.line.width = Pt(1.0)
    pr.shadow.inherit = False
_, tf = _txt(s, jx, Inches(3.52), jw, Inches(0.3))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
tf.margin_left = Inches(0.28)
_run(p, "4\u00d7 PoE Ethernet  \u2014  2 in use, 2 free", size=11, color=MUTED)

# inner application box
box(s, Emu(int(jx) + int(Inches(0.28))), Inches(4.0), Inches(4.0), Inches(1.3),
    "Multi-camera application", "Per-camera inference + composited output",
    fill=PANEL, border=VIOLET, tsize=14, ssize=11)

# --- Jetson -> laptop ---
arrow(s, Inches(8.42), Inches(3.7), Inches(0.85), Inches(0.3))
_, tf = _txt(s, Inches(8.3), Inches(3.25), Inches(1.2), Inches(0.4))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_run(p, "HTTP MJPEG", size=11, color=ACCENT, bold=True)

# --- Right: laptop browser with multi-camera tiles ---
lx, ly, lw, lh = Inches(9.45), Inches(2.65), Inches(3.3), Inches(2.55)
lap = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, ly, lw, lh)
lap.fill.solid(); lap.fill.fore_color.rgb = PANEL
lap.line.color.rgb = ACCENT; lap.line.width = Pt(1.25); lap.shadow.inherit = False
_, tf = _txt(s, lx, Inches(2.78), lw, Inches(0.4))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_run(p, "Laptop browser", size=14, color=WHITE, bold=True)
# 2x2 tiles (2 lit, 2 dim)
tile_w, tile_h = Inches(1.3), Inches(0.62)
tpos = [(Inches(9.7), Inches(3.35)), (Inches(11.15), Inches(3.35)),
        (Inches(9.7), Inches(4.12)), (Inches(11.15), Inches(4.12))]
for i, (tx, ty) in enumerate(tpos):
    lit = i < 2
    tile = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tx, ty, tile_w, tile_h)
    tile.fill.solid(); tile.fill.fore_color.rgb = ACCENT_D if lit else BG
    tile.line.color.rgb = ACCENT if lit else LINE; tile.line.width = Pt(1.0)
    tile.shadow.inherit = False
_, tf = _txt(s, lx, Inches(4.85), lw, Inches(0.3))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_run(p, "Composited multi-camera view", size=10, color=MUTED)

# caption
_, tf = _txt(s, Inches(0.7), Inches(5.95), Inches(12), Inches(1.0))
_run(tf.paragraphs[0],
     "Tested on one Jetson with 4 PoE ports: two cameras connected and the multi-camera "
     "application ran successfully — the processed stream was viewable in a laptop browser over HTTP MJPEG.",
     size=14, color=MUTED, italic=True)

# ============================================================================
# SLIDE 10 — OBSERVABILITY
# ============================================================================
s = content_slide("Observability", "Always reflects reality in the field")
bullets(s, [
    "Heartbeats every N seconds: status, active containers, and GPU telemetry.",
    "Offline watchdog flags devices automatically when heartbeats stop.",
    "Per-device deployment status, reconciled from agent ACKs.",
    "Full audit trail — every command, who issued it, when, and the device response.",
    "Event log surfaces failures (image pull, mount missing, GPU/runtime errors) with context.",
], top=Inches(2.05))

# ============================================================================
# SLIDE 11 — SECURITY
# ============================================================================
s = content_slide("Security & trust", "Touching real hardware, safely")
bullets(s, [
    "Single sign-on via Tapis OAuth2 (institutional login) — no separate passwords.",
    "Role-aware UI: admins vs. operators see the right controls.",
    "Per-device API keys: the agent authenticates every heartbeat and ACK.",
    "Owner-scoped data: users only see and act on their own devices and models.",
    "Outbound-only agent + signed sessions reduce the attack surface at the edge.",
], top=Inches(2.05))

# ============================================================================
# SLIDE 12 — TECH STACK
# ============================================================================
s = content_slide("Under the hood", "Technology choices")
bullets(s, [
    "Backend: Python + FastAPI (REST API + server-rendered UI).",
    "Messaging: MQTT for commands, HTTPS for heartbeats / ACKs.",
    "Storage: SQLite for local dev \u2192 PostgreSQL in production (same code, env-switched).",
    "Auth: Tapis OAuth2 (authorization-code flow).",
    "Edge: Python systemd agent + Docker + NVIDIA container runtime on Jetson.",
    "Packaged as a Docker image, deployed on Tapis Pods (Kubernetes).",
], top=Inches(2.05), size=18)

# ============================================================================
# SLIDE 13 — ROADMAP
# ============================================================================
s = content_slide("Roadmap", "Where we'd love NVIDIA's input")
bullets(s, [
    "DeepStream / Triton pipelines as first-class model-card runtimes.",
    "TensorRT optimization service — convert/validate engines per Jetson generation.",
    "Larger fleets + richer alerting; Orin / IGX support and MIG-aware scheduling.",
    "Self-service device onboarding and over-the-air agent updates.",
], top=Inches(2.05))

# ============================================================================
# SLIDE 14 — DISCUSSION / CLOSING
# ============================================================================
s = prs.slides.add_slide(BLANK)
_set_bg(s)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), EMU_H)
band.fill.solid(); band.fill.fore_color.rgb = ACCENT; band.line.fill.background()
_, tf = _txt(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.4))
_run(tf.paragraphs[0], "One dashboard. A fleet of Jetsons.", size=36, color=WHITE, bold=True)
p = tf.add_paragraph(); _run(p, "AI at the edge — deployed, observed, and accelerated.", size=24, color=ACCENT, bold=True)
_, tf = _txt(s, Inches(0.92), Inches(4.3), Inches(11), Inches(1.6))
_run(tf.paragraphs[0], "Discussion topics:", size=16, color=WHITE, bold=True)
for t in ["DeepStream / Triton integration", "TensorRT optimization workflows",
          "Reference Jetson generations & scaling targets"]:
    p = tf.add_paragraph(); p.space_before = Pt(4)
    _run(p, "\u2022  " + t, size=15, color=MUTED)
_, tf = _txt(s, Inches(0.92), Inches(6.2), Inches(11), Inches(0.6))
_run(tf.paragraphs[0], "Thank you  ·  Questions?  ·  Live demo available", size=16, color=WHITE)

# --- Save --------------------------------------------------------------------
out = Path(__file__).resolve().parent / "ICICLE_Edge_Platform_NVIDIA.pptx"
prs.save(str(out))
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
